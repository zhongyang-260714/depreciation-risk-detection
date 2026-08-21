#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Generate a defense PPT for the "科创企业资产折旧风险识别" competition report v4.0.
Design: 16:9, eye-friendly light gray-blue background, Microsoft YaHei, minimal text, large visuals.
"""
import json
import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

# ------------------------------------------------------------------ paths
OUT = r"D:\科创企业资产折旧算法\折旧错配：AI泡沫下科创企业资产减值风险的识别与预警_答辩PPT.pptx"
FIG_DIR = r"D:\depreciation-risk-detection\assets\report_figs_v4"
FIG_MAP = {
    "3-1": "fc8a8096ba6b44c9fc21f9a577f02387e3b883c9.png",
    "4-1": "f8175bc5f074392fbae3ce4a1692ecfabd57c576.png",
    "4-2": "e1e7af4175bced88df63bdd7a20f495e74eda5b3.png",
    "4-3": "edf020db7b4556af641d58ccb24e7f86c17250da.png",
    "4-4": "fffb4ab9096898660116648a0ff1b177d1e256ed.png",
    "4-5": "7bf8a3fb21028c9113e930971755af987566ad21.png",
    "F-1": "f21da8bdb8cd98b51419a9f7cb1368c047d8a30d.png",
    "F-2": "7526d6377b97e72f6189083c9acf1c6548cfc10d.png",
}

# also persist the corrected map
Path(r"D:/depreciation-risk-detection/assets/fig_map_v4_corrected.json").write_text(
    json.dumps(FIG_MAP, ensure_ascii=False, indent=2), encoding="utf-8")

# ------------------------------------------------------------------ design tokens
W, H = Inches(13.333), Inches(7.5)  # 16:9
BG = RGBColor(0xF5, 0xF8, 0xFB)          # light gray-blue
DARK = RGBColor(0x1A, 0x1A, 0x1A)        # near-black
NAVY = RGBColor(0x1F, 0x3A, 0x5F)        # deep blue
RED = RGBColor(0xC0, 0x39, 0x2B)         # high risk
GREEN = RGBColor(0x27, 0xAE, 0x60)       # low risk
AMBER = RGBColor(0xE0, 0xA8, 0x00)       # highlight
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "微软雅黑"

# ------------------------------------------------------------------ helpers
prs = Presentation()
prs.slide_width = W
prs.slide_height = H


def add_bg(slide, color=BG):
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    background.fill.solid()
    background.fill.fore_color.rgb = color
    background.line.fill.background()
    # send to back
    spTree = slide.shapes._spTree
    sp = background._element
    spTree.remove(sp)
    spTree.insert(2, sp)


def add_title_bar(slide, title, subtitle=None):
    # navy top bar with title
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(1.05))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    # indent slightly
    p.space_before = Pt(0); p.space_after = Pt(0)
    tf.margin_left = Inches(0.45)
    tf.margin_top = Inches(0.25)
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.name = FONT; p2.font.size = Pt(14); p2.font.color.rgb = RGBColor(0xCC,0xD6,0xE5)
        p2.alignment = PP_ALIGN.LEFT
    return bar


def add_textbox(slide, left, top, width, height, text, font_size=20,
                color=DARK, bold=False, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
                line_space=1.3):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT; p.font.size = Pt(font_size); p.font.color.rgb = color; p.font.bold = bold
    p.alignment = align
    p.line_rule = None
    # python-pptx doesn't expose line spacing multiplier easily; set space after
    p.space_after = Pt(font_size * 0.35)
    return box


def add_bullet_list(slide, left, top, width, height, items, font_size=20, color=DARK, bold_lead=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.clear()
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + item
        p.font.name = FONT; p.font.size = Pt(font_size); p.font.color.rgb = color
        p.space_after = Pt(font_size * 0.45)
        p.level = 0
    return box


def add_figure(slide, fig_key, left, top, width, height, caption=None):
    path = os.path.join(FIG_DIR, FIG_MAP[fig_key])
    if os.path.exists(path):
        pic = slide.shapes.add_picture(path, left, top, width=width)
        # keep aspect ratio; adjust top if height given? add_picture ignores height when width given.
        if caption:
            cap = slide.shapes.add_textbox(left, top + pic.height + Pt(6), width, Inches(0.35))
            tf = cap.text_frame; p = tf.paragraphs[0]
            p.text = caption; p.font.name = FONT; p.font.size = Pt(14); p.font.color.rgb = NAVY
            p.alignment = PP_ALIGN.CENTER
    else:
        add_textbox(slide, left, top, width, height, f"[图 {fig_key} 缺失]", font_size=16, color=RED)


def add_footer(slide, page_no):
    foot = slide.shapes.add_textbox(Inches(0.4), H - Inches(0.45), Inches(1.2), Inches(0.3))
    p = foot.text_frame.paragraphs[0]
    p.text = f"{page_no}"
    p.font.name = FONT; p.font.size = Pt(14); p.font.color.rgb = NAVY


def add_divider_line(slide, top):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), top, W - Inches(0.8), Pt(1.5))
    line.fill.solid(); line.fill.fore_color.rgb = RGBColor(0xCC,0xD6,0xE5); line.line.fill.background()

# ------------------------------------------------------------------ slide 1: cover
def slide_cover():
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(s)
    # decorative navy block
    blk = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.35), H)
    blk.fill.solid(); blk.fill.fore_color.rgb = NAVY; blk.line.fill.background()
    # title
    add_textbox(s, Inches(0.9), Inches(1.4), Inches(11.5), Inches(1.4),
                "折旧错配：AI 泡沫下科创企业资产减值风险的识别与预警",
                font_size=44, color=NAVY, bold=True)
    add_textbox(s, Inches(0.9), Inches(2.8), Inches(11.5), Inches(0.6),
                "基于美国十大科技公司 10-K 年报面板样本的证据链标注与五维评分体系",
                font_size=22, color=DARK)
    add_divider_line(s, Inches(3.7))
    # meta boxes
    add_textbox(s, Inches(0.9), Inches(4.1), Inches(5.5), Inches(0.5),
                "赛题编号：XH-202626", font_size=20, color=NAVY, bold=True)
    add_textbox(s, Inches(0.9), Inches(4.7), Inches(5.5), Inches(0.5),
                "第十九届“挑战杯”全国大学生课外学术科技作品竞赛 · 揭榜挂帅", font_size=18, color=DARK)
    add_textbox(s, Inches(0.9), Inches(5.8), Inches(5.5), Inches(0.5),
                "团队：待填  |  指导教师：待填", font_size=18, color=DARK)
    # big visual hint (no image, just a stylized chart glyph using shapes)
    # right side: 5D-DRS glyph
    cx, cy = Inches(10.2), Inches(5.2)
    r = Inches(1.35)
    for i, lab in enumerate(["D1", "D2", "D3", "D4", "D5"]):
        ang = 90 - i * 72
        import math
        x = cx + r * math.cos(math.radians(ang)) * 0.85 - Inches(0.28)
        y = cy - r * math.sin(math.radians(ang)) * 0.85 - Inches(0.18)
        add_textbox(s, x, y, Inches(0.6), Inches(0.4), lab, font_size=18, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    # pentagon outline
    pts = []
    for i in range(5):
        ang = math.radians(90 - i * 72)
        pts.append((cx + r*math.cos(ang), cy - r*math.sin(ang)))
    for i in range(5):
        l = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, pts[i][0], pts[i][1], pts[(i+1)%5][0], pts[(i+1)%5][1])
        l.line.color.rgb = NAVY; l.line.width = Pt(2)
    add_footer(s, 1)

# ------------------------------------------------------------------ content slides
def slide_topic():
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    add_title_bar(s, "选题与意义")
    add_textbox(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.7),
                "AI 算力资产快速迭代，折旧年限“账面”与“技术”错配，可能系统性放大科创企业资产减值风险",
                font_size=24, color=NAVY, bold=True)
    # three columns
    col_w = Inches(3.9)
    cols = [
        ("市场背景", "全球科技企业 AI 资本开支\n2024 年预计超 1,760 亿美元\n折旧会计估计是利润调节的核心杠杆之一"),
        ("现实痛点", "平安科技保险课题组调研显示\n科创企业固定资产减值识别\n缺乏统一、可验证、可复算的方法"),
        ("研究价值", "将 SEC 10-K 文本转化为\n结构化证据链与五维风险评分\n为企业、保险、监管提供共同语言"),
    ]
    for i, (ttl, body) in enumerate(cols):
        x = Inches(0.45) + i * (col_w + Inches(0.15))
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.3), col_w, Inches(3.8))
        box.fill.solid(); box.fill.fore_color.rgb = WHITE; box.line.color.rgb = RGBColor(0xCC,0xD6,0xE5)
        add_textbox(s, x + Inches(0.2), Inches(2.5), col_w - Inches(0.4), Inches(0.6),
                    ttl, font_size=22, color=NAVY, bold=True)
        add_textbox(s, x + Inches(0.2), Inches(3.2), col_w - Inches(0.4), Inches(2.6),
                    body, font_size=18, color=DARK)
    # big number
    add_textbox(s, Inches(9.8), Inches(6.5), Inches(3), Inches(0.6),
                "1,760 亿美元", font_size=32, color=RED, bold=True, align=PP_ALIGN.RIGHT)
    add_footer(s, 2)


def slide_three_qs():
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    add_title_bar(s, "核心三问")
    qs = [
        ("Q1 识别", "如何从 10-K 海量附注中\n自动定位折旧会计估计变更？"),
        ("Q2 评估", "如何把“年限延长 1 年”\n转化为可比较的风险分数？"),
        ("Q3 管理", "企业、保险、监管\n分别该做什么？"),
    ]
    for i, (q, a) in enumerate(qs):
        y = Inches(1.5) + i * Inches(1.95)
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), y, Inches(1.1), Inches(1.1))
        circ.fill.solid(); circ.fill.fore_color.rgb = NAVY; circ.line.fill.background()
        add_textbox(s, Inches(0.6), y + Inches(0.28), Inches(1.1), Inches(0.6),
                    q, font_size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(s, Inches(2.0), y + Inches(0.18), Inches(10.5), Inches(0.7),
                    a, font_size=24, color=DARK)
    add_footer(s, 3)


def slide_method():
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    add_title_bar(s, "方法论总览：四阶段技术路线")
    add_figure(s, "3-1", Inches(0.7), Inches(1.4), Inches(12), Inches(5.4),
               caption="图 3-1 研究技术路线图")
    add_footer(s, 4)


def slide_evidence():
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    add_title_bar(s, "证据链标注：把年报文本变成可验证的证据")
    # four columns
    headers = ["原文摘录", "行号/位置", "会计含义", "推断链"]
    widths = [Inches(3.2), Inches(2.4), Inches(3.3), Inches(3.2)]
    x = Inches(0.4)
    y = Inches(1.4)
    for i, h in enumerate(headers):
        hdr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, widths[i], Inches(0.55))
        hdr.fill.solid(); hdr.fill.fore_color.rgb = NAVY; hdr.line.fill.background()
        add_textbox(s, x, y + Inches(0.08), widths[i], Inches(0.4), h,
                    font_size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        x += widths[i] + Inches(0.05)
    # oracle example row
    row = [
        '"the estimated useful life of our servers and storage devices from five years to six years"',
        "Oracle FY2025 10-K\n（行号团队亲验）",
        "服务器折旧年限\n5 → 6 年",
        "折旧费用↓ $733M\n净利润↑ $573M\nEPS↑ $0.21",
    ]
    x = Inches(0.4); y = Inches(2.1)
    for i, cell in enumerate(row):
        box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, widths[i], Inches(1.6))
        box.fill.solid(); box.fill.fore_color.rgb = WHITE; box.line.color.rgb = RGBColor(0xCC,0xD6,0xE5)
        add_textbox(s, x + Inches(0.08), y + Inches(0.12), widths[i] - Inches(0.16), Inches(1.4),
                    cell, font_size=15, color=DARK)
        x += widths[i] + Inches(0.05)
    # bottom bullet
    add_bullet_list(s, Inches(0.6), Inches(4.0), Inches(12), Inches(1.8),
                    ["四列结构 = 可逐字核验：任何评委都能按图索骥回到 SEC 原文",
                     "XBRL 噪声排除 + 人工二次确认，确保信号非幻觉",
                     "30 份 10-K 年报金标准标注库，全部经作者双签确认"],
                    font_size=20)
    add_footer(s, 5)


def slide_5d_model():
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    add_title_bar(s, "5D-DRS 五维评分模型")
    # weights table
    rows = [
        ("维度", "含义", "权重", "核心输入"),
        ("D1 年限错配", "账面年限 vs 技术迭代年限", "25%", "服务器折旧年限、技术替代周期"),
        ("D2 政策保守性", "会计政策偏向激进/稳健", "20%", "折旧方法、残值率、变更历史"),
        ("D3 减值触发", "减值迹象与计提充分性", "20%", "资产减值损失、商誉减值、存货跌价"),
        ("D4 CAPEX 强度", "资本开支占收入比重", "20%", "capex_to_revenue、PP&E 净值"),
        ("D5 竞争替代", "技术被替代风险", "15%", "AI 迭代事件、产品周期、市占率"),
    ]
    # draw table manually
    col_x = [Inches(0.6), Inches(3.0), Inches(5.3), Inches(6.8)]
    col_w = [Inches(2.3), Inches(2.2), Inches(1.35), Inches(4.8)]
    row_h = Inches(0.65)
    y0 = Inches(1.5)
    for r_idx, row in enumerate(rows):
        y = y0 + r_idx * row_h
        for c_idx, cell in enumerate(row):
            rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, col_x[c_idx], y, col_w[c_idx], row_h)
            if r_idx == 0:
                rect.fill.solid(); rect.fill.fore_color.rgb = NAVY; rect.line.fill.background()
                color = WHITE; bold = True
            else:
                rect.fill.solid(); rect.fill.fore_color.rgb = WHITE; rect.line.color.rgb = RGBColor(0xCC,0xD6,0xE5)
                color = DARK; bold = (c_idx == 0)
            add_textbox(s, col_x[c_idx] + Inches(0.08), y + Inches(0.12), col_w[c_idx] - Inches(0.16), row_h - Inches(0.2),
                        cell, font_size=17, color=color, bold=bold, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_textbox(s, Inches(0.6), Inches(5.6), Inches(12), Inches(0.8),
                "综合评分 = Σ 维度得分 × 权重  →  1–5 分，4.0 以上为高风险区",
                font_size=22, color=NAVY, bold=True)
    add_footer(s, 6)


def slide_data():
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    add_title_bar(s, "数据底座：SEC 一手源 + 30 份金标准标注")
    stats = [
        ("30", "份 10-K 年报面板样本\n10 家公司 × 3 财年"),
        ("55", "个全量数值特征\n经特征工程筛选"),
        ("100%", "人工双签确认\n每条信号可回源"),
        ("0", "第三方二手平台\n未经核验的数据"),
    ]
    for i, (num, txt) in enumerate(stats):
        x = Inches(0.6) + i * Inches(3.15)
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), Inches(2.8), Inches(3.0))
        box.fill.solid(); box.fill.fore_color.rgb = WHITE; box.line.color.rgb = RGBColor(0xCC,0xD6,0xE5)
        add_textbox(s, x, Inches(2.0), Inches(2.8), Inches(0.9), num,
                    font_size=44, color=NAVY if i < 3 else RED, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(s, x + Inches(0.15), Inches(2.95), Inches(2.5), Inches(1.6), txt,
                    font_size=18, color=DARK, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.6), Inches(5.3), Inches(12), Inches(0.7),
                "数据来源：SEC EDGAR 10-K 原文；噪声排除：XBRL 标准化字段；处理：关键词矩阵 + 人工复核",
                font_size=20, color=DARK)
    add_footer(s, 7)


def slide_finding1():
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    add_title_bar(s, "发现①：风险 = 假设激进度 × 资产敞口")
    add_figure(s, "4-1", Inches(0.4), Inches(1.35), Inches(6.3), Inches(3.0), caption="图 4-1 美国十大科技公司折旧风险评分总榜")
    add_figure(s, "4-3", Inches(6.9), Inches(1.35), Inches(6.0), Inches(3.0), caption="图 4-3 CAPEX 强度与风险评分正相关")
    add_textbox(s, Inches(0.5), Inches(5.0), Inches(12), Inches(0.6),
                "制造型 / 运营型霸榜前五；轻资产设计型 / 软件型居尾——风险高低由资产结构内生决定",
                font_size=22, color=NAVY, bold=True)
    add_footer(s, 8)


def slide_finding2():
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    add_title_bar(s, "发现②：共振 + 连续延长者")
    add_figure(s, "4-4", Inches(0.4), Inches(1.35), Inches(6.3), Inches(3.0), caption="图 4-4 云厂商服务器折旧年限变化轨迹")
    add_figure(s, "4-5", Inches(6.9), Inches(1.35), Inches(6.0), Inches(3.0), caption="图 4-5 30 样本综合折旧风险评分时间序列")
    add_textbox(s, Inches(0.5), Inches(5.0), Inches(12), Inches(0.6),
                "2023–2024 年云厂商集中延长服务器折旧年限，形成行业级“折旧共振”——减值风险被集体递延",
                font_size=22, color=RED, bold=True)
    add_footer(s, 9)


def slide_finding3():
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    add_title_bar(s, "发现③④：折旧指纹与四通道管理")
    add_figure(s, "4-2", Inches(0.5), Inches(1.35), Inches(6.0), Inches(4.5), caption="图 4-2 三类资产结构的 risk profile 雷达对比")
    add_bullet_list(s, Inches(7.0), Inches(1.6), Inches(5.8), Inches(4.0),
                    ["制造型：D4 CAPEX 强度拉满，\n一旦技术路线错误，减值集中爆发",
                     "运营型：D2 政策保守性低、\nD5 竞争替代敏感",
                     "设计型：整体风险低，但需\n关注 D3 减值触发事件",
                     "管理方案：R1 审计问询 → R2 保险定价 → R3 监管预警 → R4 主动披露"],
                    font_size=20)
    add_footer(s, 10)


def slide_amazon():
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    add_title_bar(s, "反事实对照：Amazon 反向缩回 6 → 5 年")
    add_textbox(s, Inches(0.5), Inches(1.35), Inches(12.3), Inches(0.7),
                "当行业普遍延长折旧年限时，Amazon 选择缩回 1 年——这不是“不得不延”，而是管理层的真实选择",
                font_size=24, color=NAVY, bold=True)
    add_bullet_list(s, Inches(0.7), Inches(2.3), Inches(12), Inches(3.5),
                    ["2025 年亚马逊将 AWS 服务器折旧年限从 6 年缩回 5 年",
                     "2025 年经营利润因此减少约 70 亿美元，Q4'24 一次性加速折旧 9.2 亿美元",
                     "反事实意义：若延至 6 年，当期利润会被显著夸大；证明“延长”并非技术必然，而是会计估计选择",
                     "5D-DRS 模型能区分“正向延长”与“反向缩回”，避免把所有延长都判为高风险"],
                    font_size=22)
    add_footer(s, 11)


def slide_ai():
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    add_title_bar(s, "AI 技术应用：四层架构 + 可解释验证")
    add_figure(s, "F-1", Inches(0.4), Inches(1.35), Inches(6.2), Inches(3.0), caption="图 F-1 LOGO 交叉验证：预测值 vs 实际值")
    add_figure(s, "F-2", Inches(6.8), Inches(1.35), Inches(6.1), Inches(3.0), caption="图 F-2 SHAP 全局特征重要性 Top 20")
    add_bullet_list(s, Inches(0.6), Inches(5.0), Inches(12), Inches(1.5),
                    ["NLP 关键词矩阵 + DeepSeek 辅助标注 → 30 份金标准",
                     "XGBoost PoC（LOGO 交叉验证，MAE 0.418，较基线改善 ~47%）验证人工评分可被模型复现",
                     "Streamlit 智能识别系统：输入指标 → 实时评分 + SHAP 驱动因素解释"],
                    font_size=20)
    add_footer(s, 12)


def slide_solution():
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    add_title_bar(s, "落地方案：三类主体 + 四级响应")
    cols = [
        ("企业", "主动披露\n折旧假设敏感性",
         "R1\n审计问询", "R2\n保险定价"),
        ("保险公司", "将 5D-DRS 评分\n纳入科技保险费率",
         "R3\n监管预警", "R4\n主动披露"),
        ("监管机构", "监测行业级\n折旧共振信号",
         "阈值触发", "问询函"),
    ]
    col_w = Inches(3.9)
    for i, (who, act, r1, r2) in enumerate(cols):
        x = Inches(0.45) + i * (col_w + Inches(0.15))
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.5), col_w, Inches(5.0))
        box.fill.solid(); box.fill.fore_color.rgb = WHITE; box.line.color.rgb = RGBColor(0xCC,0xD6,0xE5)
        add_textbox(s, x, Inches(1.7), col_w, Inches(0.7), who, font_size=26, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(s, x + Inches(0.2), Inches(2.5), col_w - Inches(0.4), Inches(1.0), act, font_size=19, color=DARK, align=PP_ALIGN.CENTER)
        # response chips
        chip_h = Inches(0.65)
        for j, (chip, col) in enumerate([(r1, AMBER), (r2, GREEN)]):
            cy = Inches(4.0) + j * Inches(1.0)
            chipbox = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.8), cy, col_w - Inches(1.6), chip_h)
            chipbox.fill.solid(); chipbox.fill.fore_color.rgb = col; chipbox.line.fill.background()
            add_textbox(s, x + Inches(0.8), cy + Inches(0.12), col_w - Inches(1.6), chip_h, chip,
                        font_size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(s, 13)


def slide_china():
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    add_title_bar(s, "中国镜鉴：A 股三家科创企业验证")
    companies = [
        ("寒武纪", "AI 芯片设计型", "轻资产、高研发\n风险集中于 D3 减值触发"),
        ("海光信息", "CPU/DCU 制造型", "重资产、高 CAPEX\nD4 与 D1 双高"),
        ("中芯国际", "晶圆代工制造型", "5→8 年设备折旧\nD1/D4/D5 共振"),
    ]
    for i, (name, typ, risk) in enumerate(companies):
        x = Inches(0.5) + i * Inches(4.2)
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.6), Inches(3.9), Inches(4.3))
        box.fill.solid(); box.fill.fore_color.rgb = WHITE; box.line.color.rgb = RGBColor(0xCC,0xD6,0xE5)
        add_textbox(s, x, Inches(1.8), Inches(3.9), Inches(0.8), name, font_size=30, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(s, x + Inches(0.15), Inches(2.7), Inches(3.6), Inches(0.6), typ, font_size=20, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(s, x + Inches(0.2), Inches(3.5), Inches(3.5), Inches(1.8), risk, font_size=19, color=DARK, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.5), Inches(6.2), Inches(12), Inches(0.5),
                "结论：美国框架可直接迁移到中国科创板，三类资产结构的风险 profile 具有一致性",
                font_size=22, color=NAVY, bold=True)
    add_footer(s, 14)


def slide_limitations():
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    add_title_bar(s, "局限与边界：诚实声明")
    add_bullet_list(s, Inches(0.8), Inches(1.6), Inches(12), Inches(4.0),
                    ["有效独立样本量 ≈ 10（10 家公司 × 3 年面板，同公司高度相关）",
                     "XGBoost PoC 为可行性验证，不构成对新样本的预测能力声明",
                     "评分体系依赖人工确认的信号质量，未来可扩展至更大样本与半监督学习",
                     "所有结论仅限“折旧错配”风险识别，不替代审计判断或投资决策"],
                    font_size=22)
    add_textbox(s, Inches(0.8), Inches(5.8), Inches(12), Inches(0.7),
                "诚实声明不是减分项，而是让评委看到方法论边界清晰、不自欺欺人",
                font_size=22, color=GREEN, bold=True)
    add_footer(s, 15)


def slide_team():
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    add_title_bar(s, "团队与分工")
    # placeholder table
    rows = [("姓名", "专业/年级", "负责模块"), ("待填", "待填", "文献调研、方法论设计"),
            ("待填", "待填", "数据标注、10-K 证据链核验"), ("待填", "待填", "系统开发、模型工程化"),
            ("待填", "待填", "报告撰写、可视化与 PPT")]
    col_x = [Inches(0.7), Inches(3.3), Inches(5.7)]
    col_w = [Inches(2.4), Inches(2.2), Inches(6.8)]
    row_h = Inches(0.72)
    y0 = Inches(1.6)
    for r_idx, row in enumerate(rows):
        y = y0 + r_idx * row_h
        for c_idx, cell in enumerate(row):
            rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, col_x[c_idx], y, col_w[c_idx], row_h)
            if r_idx == 0:
                rect.fill.solid(); rect.fill.fore_color.rgb = NAVY; rect.line.fill.background(); color = WHITE; bold = True
            else:
                rect.fill.solid(); rect.fill.fore_color.rgb = WHITE; rect.line.color.rgb = RGBColor(0xCC,0xD6,0xE5); color = DARK; bold = False
            add_textbox(s, col_x[c_idx] + Inches(0.1), y + Inches(0.15), col_w[c_idx] - Inches(0.2), row_h - Inches(0.25),
                        cell, font_size=20, color=color, bold=bold, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_textbox(s, Inches(0.7), Inches(6.1), Inches(12), Inches(0.5),
                "提示：提交前务必将本页替换为真实姓名、专业与分工",
                font_size=18, color=RED, bold=True)
    add_footer(s, 16)


def slide_conclusion():
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    add_title_bar(s, "结论与致谢")
    add_textbox(s, Inches(0.7), Inches(1.6), Inches(12), Inches(0.8),
                "我们构建了一套“证据链标注 + 5D-DRS 评分 + 可解释 AI 验证”的折旧风险识别体系",
                font_size=30, color=NAVY, bold=True)
    add_bullet_list(s, Inches(1.0), Inches(2.7), Inches(11.5), Inches(3.0),
                    ["原创：四列证据链 + 五维评分，把文本证据变成可复算分数",
                     "可验证：30 份 SEC 一手源标注、实时评分引擎、开源代码",
                     "可落地：企业 / 保险 / 监管三类主体的四级管理方案",
                     "可迁移：美国框架在中国科创板三家样本上验证成立"],
                    font_size=24)
    add_textbox(s, Inches(0.7), Inches(6.0), Inches(12), Inches(0.7),
                "谢谢评委老师！欢迎提问",
                font_size=36, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_footer(s, 17)

# ------------------------------------------------------------------ build
slide_cover()
slide_topic()
slide_three_qs()
slide_method()
slide_evidence()
slide_5d_model()
slide_data()
slide_finding1()
slide_finding2()
slide_finding3()
slide_amazon()
slide_ai()
slide_solution()
slide_china()
slide_limitations()
slide_team()
slide_conclusion()

prs.save(OUT)
print("SAVED:", OUT)
