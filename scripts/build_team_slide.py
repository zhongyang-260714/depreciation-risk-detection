from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy, re, os

SRC = r"C:\Users\Lenovo\Desktop\折旧错配_答辩PPT.pptx"
NEW = r"C:\Users\Lenovo\Desktop\折旧错配：AI泡沫下科创企业资产减值风险的识别与预警_答辩PPT.pptx"

GOLD   = RGBColor(0xC9,0xA9,0x61)
NAVY   = RGBColor(0x24,0x3B,0x6E)
DARK   = RGBColor(0x0E,0x1B,0x2E)
RED    = RGBColor(0xE8,0x5D,0x5D)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
GRAYTX = RGBColor(0x55,0x55,0x55)
LGRAY  = RGBColor(0xE8,0xE8,0xE3)
BG     = RGBColor(0xF5,0xF8,0xFB)

prs = Presentation(SRC)
SW, SH = prs.slide_width, prs.slide_height  # EMU

def IN(v): return Inches(v)

# ---- build team slide (append, then reorder before Q&A) ----
layout = prs.slides[10].slide_layout  # content slide layout (slide 11)
ts = prs.slides.add_slide(layout)

def add_rect(slide, l,t,w,h, fill, line=None, line_w=1.0):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, IN(l),IN(t),IN(w),IN(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp

def add_text(slide, l,t,w,h, text, size=18, color=GRAYTX, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="微软雅黑"):
    tb = slide.shapes.add_textbox(IN(l),IN(t),IN(w),IN(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = font
    return tb

# top & bottom gold bars
add_rect(ts, 0, 0, 13.33, 0.06, GOLD)
add_rect(ts, 0, 7.44, 13.33, 0.06, GOLD)
# header gold square + title + tag
add_rect(ts, 0.5, 0.35, 0.7, 0.55, GOLD)
add_text(ts, 0.5, 0.35, 0.7, 0.55, "团", size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(ts, 1.4, 0.35, 8.0, 0.55, "研究团队介绍", size=26, color=NAVY, bold=True, anchor=MSO_ANCHOR.MIDDLE)
add_text(ts, 9.5, 0.35, 3.3, 0.55, "折旧错配风险识别", size=13, color=GRAYTX, anchor=MSO_ANCHOR.MIDDLE)

# two member cards
cards = [
    ("康中阳", "金融数学",
     ["负责模块",
      "• 选题与制度研究：赛题解读、中美会计准则对照、文献综述",
      "• 证据链标注与人工复核：10-K 下载定位、四列证据链、30 份全部 confirmed",
      "对应成果：第一、二、六章，附录 B"]),
    ("程亚楠", "计算机科学与技术",
     ["负责模块",
      "• 数据工程与模型验证：45 指标提取、XGBoost PoC、SHAP 分析",
      "• 系统开发：Streamlit 识别系统、FastAPI 接口、冒烟测试",
      "对应成果：附录 A/C/D/F，5.2 节"]),
]
cx = [0.6, 6.9]; cw = 5.8
for i,(name,major,lines) in enumerate(cards):
    l = cx[i]; t = 1.55
    # card body
    add_rect(ts, l, t, cw, 3.7, WHITE, line=GOLD, line_w=1.25)
    # header bar
    add_rect(ts, l, t, cw, 0.62, NAVY)
    add_text(ts, l+0.2, t, cw-0.4, 0.62, "%s   ·   %s" % (name, major), size=18, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    # body bullets
    tb = ts.shapes.add_textbox(IN(l+0.25), IN(t+0.8), IN(cw-0.5), IN(2.7))
    tf = tb.text_frame; tf.word_wrap = True
    for j,ln in enumerate(lines):
        p = tf.paragraphs[0] if j==0 else tf.add_paragraph()
        r = p.add_run(); r.text = ln
        r.font.size = Pt(14 if j==0 else 13); r.font.name = "微软雅黑"
        r.font.bold = (j==0); r.font.color.rgb = (NAVY if j==0 else GRAYTX)
        p.space_after = Pt(6)
# collaboration note
add_rect(ts, 0.6, 5.45, 12.1, 1.55, BG, line=LGRAY, line_w=1.0)
add_text(ts, 0.85, 5.55, 11.6, 1.4,
         "协作机制：全项目实行「AI 草拟 + 人工复核」双签制；关键数字（评分、金额、行号）实行双人交叉核对，复核记录留痕（review_status 字段）。报告撰写与可视化由两人共同完成。",
         size=14, color=GRAYTX, anchor=MSO_ANCHOR.MIDDLE)
# placeholder footer (global renumber fixes it)
add_text(ts, 12.2, 7.02, 1.0, 0.4, "16 / 17", size=11, color=GRAYTX, align=PP_ALIGN.RIGHT)

print("team slide added. total now:", len(prs.slides))

# ---- reorder: move team (last) before Q&A (was last) ----
ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
sldIdLst = prs.slides._sldIdLst
kids = list(sldIdLst)
# swap last two (team appended at end, Q&A just before)
kids[-1], kids[-2] = kids[-2], kids[-1]
for c in list(sldIdLst): sldIdLst.remove(c)
for c in kids: sldIdLst.append(c)
print("reordered. last two swapped.")

# ---- renumber footers to i / 17 ----
for i in range(len(prs.slides)):
    sl = prs.slides[i]
    for sh in sl.shapes:
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if re.match(r'^\d+\s*/\s*1\d$', r.text.strip()):
                        r.text = "%d / 17" % (i+1)
print("footers renumbered.")

# ---- add TOC left-panel note (slide 2) ----
toc = prs.slides[1]
add_text(toc, 0.6, 4.6, 3.5, 0.5, "研究团队介绍", size=15, color=NAVY, bold=True)
add_text(toc, 0.6, 5.05, 3.5, 0.4, "研究分工与协作机制 →", size=12, color=GRAYTX)

# ---- save to new name ----
prs.save(NEW)
print("saved to:", NEW)
if os.path.exists(SRC) and SRC != NEW:
    os.remove(SRC)
    print("removed old:", SRC)
