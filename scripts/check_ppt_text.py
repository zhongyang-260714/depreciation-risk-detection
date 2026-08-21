import re
from pptx import Presentation
P = r"C:\Users\Lenovo\Desktop\折旧错配_答辩PPT.pptx"
prs = Presentation(P)
print("slides:", len(prs.slides), "size:", prs.slide_width//914400, "x", prs.slide_height//914400)
alltext = []
for i, slide in enumerate(prs.slides):
    txts = []
    for sh in slide.shapes:
        if sh.has_text_frame:
            t = sh.text_frame.text.strip()
            if t: txts.append(t)
        if sh.has_table:
            for row in sh.table.rows:
                txts.append("[tbl] " + " | ".join(c.text.strip() for c in row.cells))
    alltext.append((i+1, txts))
    # print short header
    head = txts[0][:50] if txts else "(empty)"
    print("--- slide %02d | %s" % (i+1, head))
full = "\n".join(t for _, ts in alltext for t in ts)
# check key figures
checks = {
    "5D权重 0.25/0.20/0.20/0.20/0.15": r'0\.25',
    "30份/10家×3年": r'30',
    "变更年4.25": r'4\.25',
    "非变更年3.43": r'3\.43',
    "Burry 1760亿": r'1760',
    "Oracle 26.9%": r'26\.9',
    "Meta 20.8%": r'20\.8',
    "NVIDIA 133M": r'133',
    "NVIDIA 113M": r'113',
    "GitHub": r'[Gg]it[Hh]ub',
    "Oracle 5→6 / 6年": r'6 ?年',
    "团队": r'团队',
    "pandas": r'pandas',
}
print("\n===== figure presence in PPT =====")
for label, pat in checks.items():
    print("%-30s %s" % (label, "有" if re.search(pat, full) else "无"))
# dump any occurrence of these numbers with context
print("\n===== context snippets =====")
for pat in [r'4\.25', r'3\.43', r'1760', r'26\.9', r'20\.8', r'133', r'113', r'0\.25', r'微软雅黑|雅黑']:
    for m in re.finditer(pat, full):
        s=max(0,m.start()-25); e=min(len(full),m.end()+25)
        print("[%s] ...%s..." % (pat, full[s:e].replace("\n"," ")))
        break
