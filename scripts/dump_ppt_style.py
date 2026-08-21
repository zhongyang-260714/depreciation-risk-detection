from pptx import Presentation
from pptx.util import Emu
P = r"C:\Users\Lenovo\Desktop\折旧错配_答辩PPT.pptx"
prs = Presentation(P)
SW, SH = prs.slide_width, prs.slide_height
print("slide WxH EMU:", SW, SH, "-> in:", SW/914400, SH/914400)

def dump(idx):
    print("\n========== SLIDE %d ==========" % (idx+1))
    s = prs.slides[idx]
    for sh in s.shapes:
        try:
            x,y,w,h = sh.left, sh.top, sh.width, sh.height
        except Exception:
            x=y=w=h=None
        fill=None
        try:
            if sh.fill.type is not None:
                fill = sh.fill.fore_color.rgb if sh.fill.fore_color and sh.fill.fore_color.type else None
        except Exception:
            fill=None
        txt = (sh.text_frame.text[:30] if sh.has_text_frame else ("[TABLE %dx%d"%(len(sh.table.rows),len(sh.table.columns)) if sh.has_table else ""))
        print("  %-12s L=%s T=%s W=%s H=%s fill=%s | %s" % (
            sh.shape_type, round(x/914400,2) if x is not None else '?', round(y/914400,2) if y is not None else '?',
            round(w/914400,2) if w is not None else '?', round(h/914400,2) if h is not None else '?', fill, txt.replace("\n"," ")))

for i in [1,2,11]:  # TOC(2), 智能识别系统(12)
    dump(i)
