import shutil, zipfile, re, os, sys
from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn

SRC = r"C:\Users\Lenovo\Desktop\折旧错配_答辩PPT.pptx"
TMP = r"C:\Users\Lenovo\Desktop\__tmp_reorder.pptx"
FINAL = SRC  # overwrite original after validation

# ---------- Step 1: reorder slideIdLst + patch theme (zip level) ----------
z = zipfile.ZipFile(SRC, 'r')
names = z.namelist()

# 1a. reorder presentation.xml sldIdLst
pres_xml = z.read('ppt/presentation.xml')
root = etree.fromstring(pres_xml)
ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
sldIdLst = root.find('{%s}sldIdLst' % ns_p)
children = list(sldIdLst)
print("current order count:", len(children))
# desired permutation (0-based indices): move idx13 to right after idx11
perm = [0,1,2,3,4,5,6,7,8,9,10,11,13,12,14,15]
assert len(perm) == len(children), "perm mismatch"
new_children = [children[i] for i in perm]
# rebuild sldIdLst
for c in list(sldIdLst):
    sldIdLst.remove(c)
for c in new_children:
    sldIdLst.append(c)
pres_xml2 = etree.tostring(root, xml_declaration=True, encoding='UTF-8', standalone=True)

# 1b. patch theme1.xml: ea -> 微软雅黑, Hans -> 微软雅黑
theme_xml = z.read('ppt/theme/theme1.xml').decode('utf-8')
theme_xml = theme_xml.replace('<a:ea typeface=""/>', '<a:ea typeface="微软雅黑"/>')
theme_xml = theme_xml.replace('<a:font script="Hans" typeface="宋体"/>', '<a:font script="Hans" typeface="微软雅黑"/>')
# also ensure any existing 新細明體 Hant etc unaffected; only Hans matters

# 1c. write new zip
with zipfile.ZipFile(TMP, 'w', zipfile.ZIP_DEFLATED) as zout:
    for n in names:
        data = z.read(n)
        if n == 'ppt/presentation.xml':
            zout.writestr(n, pres_xml2)
        elif n == 'ppt/theme/theme1.xml':
            zout.writestr(n, theme_xml)
        else:
            zout.writestr(n, data)
z.close()
print("temp written:", TMP)

# ---------- Step 2: set run-level east_asian_font (python-pptx) ----------
prs = Presentation(TMP)
cnt_runs = 0
cnt_set = 0
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    cnt_runs += 1
                    # set east asian font only (latin untouched -> stays Calibri)
                    try:
                        run.font.east_asian_font = "微软雅黑"
                        cnt_set += 1
                    except Exception as e:
                        pass
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            cnt_runs += 1
                            try:
                                run.font.east_asian_font = "微软雅黑"
                                cnt_set += 1
                            except Exception:
                                pass
print("runs scanned:", cnt_runs, "east_asian set:", cnt_set)
prs.save(FINAL)
print("saved final:", FINAL)

# validate
prs2 = Presentation(FINAL)
print("final slide count:", len(prs2.slides))
# verify order of chapters
import re as _re
order = []
for slide in prs2.slides:
    chap = ""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                t = "".join(r.text for r in p.runs)
                m = _re.match(r'^\s*(\d{2})\b', t.strip())
                if m:
                    chap = m.group(1); break
        if chap: break
    order.append(chap)
print("chapter order:", order)
