import sys
from pptx import Presentation
from pptx.util import Emu

PATH = r"D:\depreciation-risk-detection\..\..\..\Users\Lenovo\Desktop\折旧错配_答辩PPT.pptx"
# normalize to absolute
import os
PATH = r"C:\Users\Lenovo\Desktop\折旧错配_答辩PPT.pptx"

prs = Presentation(PATH)
print("slide count:", len(prs.slides))
print("slide size: %d x %d EMU  (%.2f x %.2f inch)" % (prs.slide_width, prs.slide_height, prs.slide_width/914400, prs.slide_height/914400))

for i, slide in enumerate(prs.slides):
    texts = []
    chap_candidate = None
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = "".join(r.text for r in para.runs)
                if t.strip():
                    texts.append(t.strip())
        # check tables
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    t = cell.text.strip()
                    if t:
                        texts.append("[tbl]"+t)
    # heuristic: first non-empty text is title
    head = texts[0] if texts else "(no text)"
    # look for chapter-like number
    import re
    chap = ""
    for t in texts:
        m = re.match(r'^\s*(\d{2})\b', t)
        if m:
            chap = m.group(1)
            break
    print("---- slide %02d | chap=%s" % (i+1, chap))
    print("     title:", head[:60])
    # print all texts truncated
    for t in texts[:6]:
        print("      ·", t[:70])
