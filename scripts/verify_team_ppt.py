import zipfile, re
from pptx import Presentation
from pptx.oxml.ns import qn
P = r"C:\Users\Lenovo\Desktop\折旧错配：AI泡沫下科创企业资产减值风险的识别与预警_答辩PPT.pptx"
prs = Presentation(P)
print("slide count:", len(prs.slides))

# display order via sldIdLst
pres = prs.slides._sldIdLst.getparent().getparent()  # presentation element
ns_p='http://schemas.openxmlformats.org/presentationml/2006/main'
sldIdLst = prs.slides._sldIdLst
order=[]
for sid in sldIdLst:
    rid=sid.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    rel=prs.slides.part.part.rels[rid]
    order.append(rel.target_ref)
print("display file order:", [o.split('/')[-1] for o in order])

for i, sl in enumerate(prs.slides, 1):
    head=""
    footer=""
    has_team=False
    for sh in sl.shapes:
        if sh.has_text_frame:
            t=sh.text_frame.text.strip()
            if t and not head:
                head=t.split("\n")[0][:40]
            if re.match(r'^\d+ / 17$', t.strip()):
                footer=t.strip()
            if "团队" in t:
                has_team=True
    print("pos%02d | %-40s | footer=%s | team=%s" % (i, head, footer, has_team))

# team slide detail
print("\n--- team slide (find '研究团队介绍') ---")
for sl in prs.slides:
    if sl.shapes and any('研究团队介绍' in (s.text_frame.text if s.has_text_frame else '') for s in sl.shapes):
        for sh in sl.shapes:
            if sh.has_text_frame:
                txt=sh.text_frame.text.strip()
                if txt: print("  |", txt.replace("\n"," / ")[:120])
        break

# zip integrity
z=zipfile.ZipFile(P)
print("\nzip testzip:", z.testzip())
# theme check still 微软雅黑
th=z.read('ppt/theme/theme1.xml').decode('utf-8')
print("ea=微软雅黑:", '<a:ea typeface="微软雅黑"/>' in th, "| Hans=微软雅黑:", '<a:font script="Hans" typeface="微软雅黑"/>' in th)
